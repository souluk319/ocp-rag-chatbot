"""문서를 적당한 크기의 청크로 쪼갬. 문단 경계를 기준으로 자르고 overlap으로 맥락 유지."""
import os
import re
from dataclasses import dataclass, field
from typing import Optional

from src.config import CHUNK_SIZE, CHUNK_OVERLAP


@dataclass
class Chunk:
    """하나의 텍스트 청크"""
    text: str
    metadata: dict = field(default_factory=dict)
    chunk_id: Optional[str] = None

    def __post_init__(self):
        if self.chunk_id is None:
            source = self.metadata.get("source", "unknown")
            idx = self.metadata.get("chunk_index", 0)
            self.chunk_id = f"{source}::{idx}"


class DocumentParser:
    """다양한 문서 포맷을 텍스트로 변환"""

    @staticmethod
    def parse_file(filepath: str) -> str:
        ext = os.path.splitext(filepath)[1].lower()
        if ext == ".txt" or ext == ".md":
            return DocumentParser._read_text(filepath)
        elif ext == ".pdf":
            return DocumentParser._read_pdf(filepath)
        elif ext == ".docx":
            return DocumentParser._read_docx(filepath)
        elif ext == ".pptx":
            return DocumentParser._read_pptx(filepath)
        else:
            # 알 수 없는 형식은 텍스트로 시도
            return DocumentParser._read_text(filepath)

    @staticmethod
    def _read_text(filepath: str) -> str:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    def _read_pdf(filepath: str) -> str:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)

    @staticmethod
    def _read_docx(filepath: str) -> str:
        from docx import Document
        doc = Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    @staticmethod
    def _read_pptx(filepath: str) -> str:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_text.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides_text)


class Chunker:
    """Sliding Window 방식으로 문서를 청킹

    전략:
    1. 문단(paragraph) 단위로 먼저 분리
    2. 문단들을 CHUNK_SIZE 이내로 결합
    3. CHUNK_OVERLAP만큼 겹치게 윈도우 슬라이드
    4. 각 청크에 소스 파일명, 섹션 제목 등 메타데이터 부착
    """

    def __init__(self, chunk_size: int = CHUNK_SIZE, chunk_overlap: int = CHUNK_OVERLAP):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    # 기술 용어 보호 리스트 (camelCase 분리하면 안 되는 단어)
    _PROTECTED_TERMS = {
        "StatefulSet", "ReplicaSet", "DaemonSet", "ConfigMap", "CronJob",
        "PersistentVolume", "PersistentVolumeClaim", "StorageClass",
        "ClusterRole", "RoleBinding", "ClusterRoleBinding",
        "ServiceAccount", "ResourceQuota", "LimitRange",
        "NetworkPolicy", "IngressController", "ImageStream",
        "BuildConfig", "DeploymentConfig", "ReplicationController",
        "CustomResource", "HorizontalPodAutoscaler", "OpenShift",
        "MachineSet", "MachineConfig", "NodePort", "ClusterIP",
        "LoadBalancer", "ReadWriteOnce", "ReadOnlyMany", "ReadWriteMany",
    }

    @staticmethod
    def _normalize_text(text: str) -> str:
        """스크래핑 데이터의 공백 누락/태그 잔여물 정리"""
        # 1. 기술 용어를 플레이스홀더로 보호
        placeholders = {}
        for i, term in enumerate(Chunker._PROTECTED_TERMS):
            placeholder = f"\x00TERM{i}\x00"
            if term in text:
                text = text.replace(term, placeholder)
                placeholders[placeholder] = term

        # 2. 소문자+대문자 경계에 공백 삽입 (붙은 단어 분리)
        # 예: "DeploymentandReplicaSet" → "Deploymentand Replica Set" → (복원 후 정상)
        text = re.sub(r"([a-z])([A-Z])", r"\1 \2", text)
        # 마침표/쉼표 뒤 공백 없는 경우 추가
        text = re.sub(r"([.,!?])([A-Z가-힣])", r"\1 \2", text)
        # 연속 공백 정리
        text = re.sub(r"[ \t]+", " ", text)

        # 3. 기술 용어 복원
        for placeholder, term in placeholders.items():
            text = text.replace(placeholder, term)

        return text

    def chunk_text(self, text: str, source: str = "unknown") -> list[Chunk]:
        """텍스트를 청크 리스트로 분할"""
        # 빈 텍스트 처리
        text = text.strip()
        if not text:
            return []
        # 텍스트 정규화 (스크래핑 아티팩트 정리)
        text = self._normalize_text(text)

        # 1단계: 문단 단위 분리
        paragraphs = self._split_paragraphs(text)

        # 2단계: 문단들을 chunk_size 이내로 결합 + overlap 슬라이딩
        chunks = self._sliding_window(paragraphs, source)
        return chunks

    def chunk_file(self, filepath: str) -> list[Chunk]:
        """파일을 파싱하고 청킹"""
        text = DocumentParser.parse_file(filepath)
        source = os.path.basename(filepath)
        return self.chunk_text(text, source=source)

    def chunk_directory(self, dirpath: str) -> list[Chunk]:
        """디렉토리 재귀 순회하면서 지원 포맷 전부 청킹"""
        all_chunks = []
        supported = {".txt", ".md", ".pdf", ".docx", ".pptx"}
        # FIXME: 대용량 PDF (100+ 페이지) 들어오면 메모리 터질 수 있음. 나중에 스트리밍 처리 필요
        for root, _, files in os.walk(dirpath):
            for fname in sorted(files):
                if fname.startswith("~$") or fname.startswith("."):
                    continue
                ext = os.path.splitext(fname)[1].lower()
                if ext in supported:
                    filepath = os.path.join(root, fname)
                    chunks = self.chunk_file(filepath)
                    all_chunks.extend(chunks)
        return all_chunks

    def _split_paragraphs(self, text: str) -> list[str]:
        """텍스트를 문단 단위로 분리 (빈 줄 기준)"""
        raw_paragraphs = re.split(r"\n\s*\n", text)
        paragraphs = []
        for p in raw_paragraphs:
            p = p.strip()
            if p:
                paragraphs.append(p)
        return paragraphs

    def _sliding_window(self, paragraphs: list[str], source: str) -> list[Chunk]:
        """문단들을 sliding window로 결합하여 청크 생성"""
        if not paragraphs:
            return []

        chunks = []
        current_texts: list[str] = []
        current_len = 0
        chunk_index = 0

        for para in paragraphs:
            para_len = len(para)

            # 단일 문단이 chunk_size보다 큰 경우 → 문장 단위로 분할
            if para_len > self.chunk_size:
                # 현재 버퍼 먼저 flush
                if current_texts:
                    chunk_text = "\n\n".join(current_texts)
                    chunks.append(Chunk(
                        text=chunk_text,
                        metadata={"source": source, "chunk_index": chunk_index},
                    ))
                    chunk_index += 1
                    # overlap: 마지막 일부 문단 유지
                    current_texts, current_len = self._apply_overlap(current_texts)

                # 큰 문단을 문장 단위로 분할
                sentences = re.split(r"(?<=[.!?。])\s+", para)
                sent_buffer = []
                sent_len = 0
                for sent in sentences:
                    if sent_len + len(sent) > self.chunk_size and sent_buffer:
                        chunk_text = " ".join(sent_buffer)
                        chunks.append(Chunk(
                            text=chunk_text,
                            metadata={"source": source, "chunk_index": chunk_index},
                        ))
                        chunk_index += 1
                        sent_buffer = []
                        sent_len = 0
                    sent_buffer.append(sent)
                    sent_len += len(sent)
                if sent_buffer:
                    remaining = " ".join(sent_buffer)
                    current_texts = [remaining]
                    current_len = len(remaining)
                continue

            # 버퍼에 추가했을 때 크기 초과 → flush
            if current_len + para_len + 2 > self.chunk_size and current_texts:
                chunk_text = "\n\n".join(current_texts)
                chunks.append(Chunk(
                    text=chunk_text,
                    metadata={"source": source, "chunk_index": chunk_index},
                ))
                chunk_index += 1
                current_texts, current_len = self._apply_overlap(current_texts)

            current_texts.append(para)
            current_len += para_len + 2  # "\n\n" 구분자

        # 남은 버퍼 flush
        if current_texts:
            chunk_text = "\n\n".join(current_texts)
            chunks.append(Chunk(
                text=chunk_text,
                metadata={"source": source, "chunk_index": chunk_index},
            ))

        return chunks

    def _apply_overlap(self, texts: list[str]) -> tuple[list[str], int]:
        """overlap 적용: 뒤에서부터 overlap 크기만큼 문단 유지"""
        if not texts or self.chunk_overlap <= 0:
            return [], 0
        overlap_texts = []
        overlap_len = 0
        for t in reversed(texts):
            if overlap_len + len(t) > self.chunk_overlap:
                break
            overlap_texts.insert(0, t)
            overlap_len += len(t) + 2
        return overlap_texts, overlap_len
