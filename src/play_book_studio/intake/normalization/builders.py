from __future__ import annotations

# captured source에서 canonical study asset을 조립하는 web/pdf 빌더 모음.

from collections.abc import Callable
from pathlib import Path
import re
from typing import Any

from play_book_studio.ingestion.models import SourceManifestEntry
from play_book_studio.ingestion.normalize import extract_sections

from ..models import CustomerPackDraftRecord
from ..planner import CustomerPackPlanner
from .pdf import (
    extract_pdf_markdown_with_docling,
    extract_pdf_markdown_with_docling_ocr,
    extract_pdf_outline,
    extract_pdf_pages,
)
from .markitdown_adapter import MARKITDOWN_SOURCE_TYPES, convert_with_markitdown
from .pdf_rows import (
    _build_pdf_rows,
    _build_pdf_rows_from_docling_markdown,
    _docling_korean_quality_ok,
)

try:  # pragma: no cover - optional runtime dependency
    from docling.document_converter import DocumentConverter
except Exception:  # noqa: BLE001
    DocumentConverter = None

try:  # pragma: no cover - optional runtime dependency
    from docx import Document as DocxDocument
    from docx.document import Document as DocxDocumentType
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph
except Exception:  # noqa: BLE001
    DocxDocument = None
    DocxDocumentType = None
    CT_Tbl = None
    CT_P = None
    DocxTable = None
    DocxParagraph = None

try:  # pragma: no cover - optional runtime dependency
    from pptx import Presentation
except Exception:  # noqa: BLE001
    Presentation = None

try:  # pragma: no cover - optional runtime dependency
    from openpyxl import load_workbook
except Exception:  # noqa: BLE001
    load_workbook = None


def build_canonical_book(
    record: CustomerPackDraftRecord,
    *,
    extract_pdf_markdown_with_docling_fn: Callable[[Path], str] = extract_pdf_markdown_with_docling,
    extract_pdf_markdown_with_docling_ocr_fn: Callable[[Path], str] = extract_pdf_markdown_with_docling_ocr,
    extract_pdf_outline_fn: Callable[[Path], list] = extract_pdf_outline,
    extract_pdf_pages_fn: Callable[[Path], list[str]] = extract_pdf_pages,
):
    if record.request.source_type in MARKITDOWN_SOURCE_TYPES:
        return _build_markitdown_first_canonical_book(
            record,
            extract_pdf_markdown_with_docling_fn=extract_pdf_markdown_with_docling_fn,
            extract_pdf_markdown_with_docling_ocr_fn=extract_pdf_markdown_with_docling_ocr_fn,
            extract_pdf_outline_fn=extract_pdf_outline_fn,
            extract_pdf_pages_fn=extract_pdf_pages_fn,
        )
    if record.request.source_type == "pdf":
        return _build_pdf_canonical_book(
            record,
            extract_pdf_markdown_with_docling_fn=extract_pdf_markdown_with_docling_fn,
            extract_pdf_markdown_with_docling_ocr_fn=extract_pdf_markdown_with_docling_ocr_fn,
            extract_pdf_outline_fn=extract_pdf_outline_fn,
            extract_pdf_pages_fn=extract_pdf_pages_fn,
        )
    if record.request.source_type == "web":
        return _build_web_canonical_book(record)
    if record.request.source_type in {"md", "asciidoc", "txt"}:
        return _build_text_canonical_book(record)
    if record.request.source_type == "docx":
        return _build_docx_canonical_book(record)
    if record.request.source_type == "pptx":
        return _build_pptx_canonical_book(record)
    if record.request.source_type == "xlsx":
        return _build_xlsx_canonical_book(record)
    if record.request.source_type == "image":
        return _build_image_canonical_book(record)
    raise ValueError("지원하지 않는 source_type입니다.")


def _append_normalization_note(book, note: str):
    book.notes = tuple(book.notes) + (note,)
    return book


def _build_markitdown_first_canonical_book(
    record: CustomerPackDraftRecord,
    *,
    extract_pdf_markdown_with_docling_fn: Callable[[Path], str],
    extract_pdf_markdown_with_docling_ocr_fn: Callable[[Path], str],
    extract_pdf_outline_fn: Callable[[Path], list],
    extract_pdf_pages_fn: Callable[[Path], list[str]],
):
    try:
        return _append_normalization_note(
            _build_markitdown_canonical_book(record),
            f"Normalization backend: MarkItDown ({record.request.source_type} -> markdown).",
        )
    except Exception as exc:
        fallback_builder = {
            "pdf": lambda: _build_pdf_canonical_book(
                record,
                extract_pdf_markdown_with_docling_fn=extract_pdf_markdown_with_docling_fn,
                extract_pdf_markdown_with_docling_ocr_fn=extract_pdf_markdown_with_docling_ocr_fn,
                extract_pdf_outline_fn=extract_pdf_outline_fn,
                extract_pdf_pages_fn=extract_pdf_pages_fn,
            ),
            "docx": lambda: _build_docx_canonical_book(record),
            "pptx": lambda: _build_pptx_canonical_book(record),
            "xlsx": lambda: _build_xlsx_canonical_book(record),
        }[record.request.source_type]
        return _append_normalization_note(
            fallback_builder(),
            f"Normalization backend: legacy fallback after MarkItDown failure ({exc.__class__.__name__}).",
        )


def _build_markitdown_canonical_book(record: CustomerPackDraftRecord):
    capture_path = Path(record.capture_artifact_path)
    if not capture_path.exists():
        raise FileNotFoundError(f"captured artifact를 찾을 수 없습니다: {capture_path}")
    markdown = convert_with_markitdown(capture_path)
    if _markitdown_output_is_low_confidence(record.request.source_type, markdown):
        raise ValueError(f"MarkItDown produced low-confidence {record.request.source_type} output.")
    normalized = _normalize_text_block(
        "md",
        markdown,
        origin_source_type=record.request.source_type,
        book_title=record.plan.title,
    )
    if not normalized:
        raise ValueError(f"{record.request.source_type.upper()}에서 canonical section을 만들지 못했습니다.")
    return _build_structured_text_canonical_book(
        record,
        source_type="md",
        normalized=normalized,
        origin_source_type=record.request.source_type,
    )


def _markitdown_output_is_low_confidence(source_type: str, markdown: str) -> bool:
    stripped = str(markdown or "").strip()
    if not stripped:
        return True
    if source_type == "pdf":
        compact = stripped.replace("\ufeff", "").strip()
        if compact.lower().startswith("%pdf-"):
            return True
    return False


def _build_web_canonical_book(record: CustomerPackDraftRecord):
    capture_path = Path(record.capture_artifact_path)
    if not capture_path.exists():
        raise FileNotFoundError(f"captured artifact를 찾을 수 없습니다: {capture_path}")

    html = capture_path.read_text(encoding="utf-8")
    manifest_entry = SourceManifestEntry(
        book_slug=record.plan.book_slug,
        title=record.plan.title,
        source_url=record.request.uri,
        viewer_path=f"/playbooks/customer-packs/{record.draft_id}/index.html",
        high_value=False,
        viewer_strategy="internal_text",
        body_language_guess=record.request.language_hint,
        approval_status="derived",
    )
    sections = extract_sections(html, manifest_entry)
    if not sections:
        raise ValueError("capture한 문서에서 canonical section을 만들지 못했습니다.")
    rows = [section.to_dict() for section in sections]
    return CustomerPackPlanner().build_canonical_book(rows, request=record.request)


def _build_pdf_canonical_book(
    record: CustomerPackDraftRecord,
    *,
    extract_pdf_markdown_with_docling_fn: Callable[[Path], str],
    extract_pdf_markdown_with_docling_ocr_fn: Callable[[Path], str],
    extract_pdf_outline_fn: Callable[[Path], list],
    extract_pdf_pages_fn: Callable[[Path], list[str]],
):
    capture_path = Path(record.capture_artifact_path)
    if not capture_path.exists():
        raise FileNotFoundError(f"captured artifact를 찾을 수 없습니다: {capture_path}")
    try:
        markdown = extract_pdf_markdown_with_docling_fn(capture_path)
        markdown_rows = _build_pdf_rows_from_docling_markdown(markdown, record)
        # Quality gate: docling sometimes merges Korean characters without spaces
        # (e.g. "컨트롤플레인정보"). Detect and fall back to pypdf in that case.
        if markdown_rows and _docling_korean_quality_ok(markdown_rows):
            return CustomerPackPlanner().build_canonical_book(markdown_rows, request=record.request)
    except Exception:
        pass
    try:
        markdown = extract_pdf_markdown_with_docling_ocr_fn(capture_path)
        markdown_rows = _build_pdf_rows_from_docling_markdown(markdown, record)
        if markdown_rows:
            return CustomerPackPlanner().build_canonical_book(markdown_rows, request=record.request)
    except Exception:
        pass
    pages = extract_pdf_pages_fn(capture_path)
    outline = extract_pdf_outline_fn(capture_path)
    rows = _build_pdf_rows(pages, record, outline=outline)
    if not rows:
        raise ValueError("PDF에서 canonical section을 만들지 못했습니다.")
    return CustomerPackPlanner().build_canonical_book(rows, request=record.request)


_MARKDOWN_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*?)\s*$")
_ASCIIDOC_HEADING_RE = re.compile(r"^(={1,6})\s+(.*?)\s*$")
_NUMERIC_HEADING_RE = re.compile(r"^(\d+(?:\.\d+){0,5})\.?\s+(.*?)\s*$")
_FENCED_CODE_START_RE = re.compile(r"^\s*(?P<fence>`{3,}|~{3,})(?P<language>[\w.+-]*)\s*$")
_TABLE_SEPARATOR_CELL_RE = re.compile(r"^:?-{3,}:?$")
_PDF_INLINE_PAGE_MARKER_RE = re.compile(r"^[가-힣A-Za-z][가-힣A-Za-z0-9\s()./_-]{1,24}\s+\d+$")


def _slug_anchor(value: str, *, ordinal: int) -> str:
    cleaned = re.sub(r"[^a-z0-9가-힣]+", "-", value.strip().lower())
    cleaned = re.sub(r"-{2,}", "-", cleaned).strip("-")
    return cleaned or f"section-{ordinal}"


def _is_fenced_code_closer(line: str, *, fence_char: str, minimum_width: int) -> bool:
    stripped = (line or "").strip()
    return bool(
        stripped
        and len(stripped) >= minimum_width
        and set(stripped) == {fence_char}
    )


def _replace_fenced_code_blocks(text: str) -> str:
    lines = text.split("\n")
    normalized_lines: list[str] = []
    index = 0
    while index < len(lines):
        line = lines[index]
        match = _FENCED_CODE_START_RE.match(line)
        if match is None:
            normalized_lines.append(line)
            index += 1
            continue
        fence = match.group("fence")
        language = str(match.group("language") or "").strip().lower()
        body_lines: list[str] = []
        index += 1
        while index < len(lines):
            candidate = lines[index]
            if _is_fenced_code_closer(candidate, fence_char=fence[0], minimum_width=len(fence)):
                index += 1
                break
            body_lines.append(candidate)
            index += 1
        attrs = f' language="{language}"' if language else ""
        body = "\n".join(body_lines).strip("\n")
        normalized_lines.append(f"[CODE{attrs}]\n{body}\n[/CODE]")
    return "\n".join(normalized_lines)


def _split_markdown_table_row(line: str) -> list[str] | None:
    stripped = (line or "").strip()
    if "|" not in stripped:
        return None
    inner = stripped[1:-1] if stripped.startswith("|") and stripped.endswith("|") else stripped
    cells = [cell.strip() for cell in inner.split("|")]
    if len(cells) < 2:
        return None
    return cells


def _is_markdown_table_separator(row: list[str]) -> bool:
    return bool(row) and all(_TABLE_SEPARATOR_CELL_RE.fullmatch(cell or "") for cell in row)


def _replace_markdown_tables(text: str) -> str:
    lines = text.split("\n")
    normalized_lines: list[str] = []
    index = 0
    while index < len(lines):
        stripped = lines[index].strip()
        if stripped.startswith("[CODE"):
            normalized_lines.append(lines[index])
            index += 1
            while index < len(lines):
                normalized_lines.append(lines[index])
                if lines[index].strip() == "[/CODE]":
                    index += 1
                    break
                index += 1
            continue
        if stripped.startswith("[TABLE"):
            normalized_lines.append(lines[index])
            index += 1
            while index < len(lines):
                normalized_lines.append(lines[index])
                if lines[index].strip() == "[/TABLE]":
                    index += 1
                    break
                index += 1
            continue
        header = _split_markdown_table_row(lines[index])
        separator = _split_markdown_table_row(lines[index + 1]) if index + 1 < len(lines) else None
        if (
            header
            and separator
            and len(header) == len(separator)
            and _is_markdown_table_separator(separator)
        ):
            table_rows = [header]
            index += 2
            while index < len(lines):
                row = _split_markdown_table_row(lines[index])
                if row is None or len(row) != len(header):
                    break
                table_rows.append(row)
                index += 1
            normalized_lines.append(_table_to_tagged_text(table_rows))
            continue
        normalized_lines.append(lines[index])
        index += 1
    return "\n".join(normalized_lines)


def _pdf_title_tokens(title: str) -> tuple[str, ...]:
    compact = re.sub(r"\([^)]*\)", " ", str(title or ""))
    compact = re.sub(r"^\d+(?:[.\-]\d+)*\s*", " ", compact).strip()
    tokens = [
        token
        for token in re.findall(r"[A-Za-z가-힣]{2,}", compact)
        if token.lower() not in {"openshift", "container", "platform"}
    ]
    return tuple(dict.fromkeys(tokens))


def _should_drop_pdf_inline_page_marker(line: str, *, book_title: str) -> bool:
    stripped = (line or "").strip()
    if not stripped or not _PDF_INLINE_PAGE_MARKER_RE.fullmatch(stripped):
        return False
    lowered = stripped.lower()
    return any(token.lower() in lowered for token in _pdf_title_tokens(book_title))


def _clean_pdf_markitdown_text(text: str, *, book_title: str) -> str:
    cleaned_lines: list[str] = []
    for raw_line in str(text or "").replace("\x00", "").splitlines():
        if _should_drop_pdf_inline_page_marker(raw_line, book_title=book_title):
            continue
        cleaned_lines.append(raw_line)
    return "\n".join(cleaned_lines)


def _looks_like_table_fragment(text: str) -> bool:
    stripped = (text or "").strip()
    return stripped.count("|") >= 2 or (stripped.startswith("|") and stripped.endswith("|"))


def _normalize_text_block(
    source_type: str,
    text: str,
    *,
    origin_source_type: str = "",
    book_title: str = "",
) -> str:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    if origin_source_type == "pdf":
        normalized = _clean_pdf_markitdown_text(normalized, book_title=book_title)
    if source_type in {"md", "asciidoc"}:
        normalized = _replace_fenced_code_blocks(normalized)
    if source_type == "md":
        normalized = _replace_markdown_tables(normalized)
    return normalized.strip()


def _text_heading_match(
    source_type: str,
    line: str,
    *,
    origin_source_type: str = "",
) -> tuple[int, str] | None:
    stripped = line.strip()
    if not stripped:
        return None
    if source_type == "md":
        match = _MARKDOWN_HEADING_RE.match(stripped)
        if match:
            return len(match.group(1)), match.group(2).strip()
    if source_type == "asciidoc":
        match = _ASCIIDOC_HEADING_RE.match(stripped)
        if match:
            return len(match.group(1)), match.group(2).strip()
    match = _NUMERIC_HEADING_RE.match(stripped)
    if match:
        if origin_source_type == "pdf" and _looks_like_table_fragment(stripped):
            return None
        return match.group(1).count(".") + 1, f"{match.group(1)} {match.group(2).strip()}".strip()
    return None


def _build_text_canonical_book(record: CustomerPackDraftRecord):
    capture_path = Path(record.capture_artifact_path)
    if not capture_path.exists():
        raise FileNotFoundError(f"captured artifact를 찾을 수 없습니다: {capture_path}")

    normalized = _normalize_text_block(record.request.source_type, capture_path.read_text(encoding="utf-8"))
    return _build_structured_text_canonical_book(
        record,
        source_type=record.request.source_type,
        normalized=normalized,
    )


def _build_structured_text_canonical_book(
    record: CustomerPackDraftRecord,
    *,
    source_type: str,
    normalized: str,
    origin_source_type: str = "",
):
    if not normalized:
        raise ValueError("텍스트 소스에서 canonical section을 만들지 못했습니다.")

    viewer_base = f"/playbooks/customer-packs/{record.draft_id}/index.html"
    current_heading = record.plan.title or record.plan.book_slug
    current_level = 1
    current_body: list[str] = []
    path_stack: list[str] = [current_heading]
    rows: list[dict[str, object]] = []

    def flush() -> None:
        ordinal = len(rows) + 1
        heading = current_heading.strip() or f"Section {ordinal}"
        anchor = _slug_anchor(heading, ordinal=ordinal)
        body = "\n".join(current_body).strip()
        if not body:
            return
        section_path = tuple(path_stack[:current_level]) if path_stack else (heading,)
        rows.append(
            {
                "book_slug": record.plan.book_slug,
                "book_title": record.plan.title,
                "heading": heading,
                "section_level": current_level,
                "section_path": list(section_path),
                "anchor": anchor,
                "source_url": record.request.uri,
                "viewer_path": f"{viewer_base}#{anchor}",
                "text": body,
            }
        )

    for line in normalized.splitlines():
        heading = _text_heading_match(
            source_type,
            line,
            origin_source_type=origin_source_type,
        )
        if heading:
            if not rows and not "\n".join(current_body).strip():
                current_level, current_heading = heading
                path_stack = path_stack[: max(current_level - 1, 0)]
                path_stack.append(current_heading)
                current_body = []
                continue
            flush()
            current_level, current_heading = heading
            path_stack = path_stack[: max(current_level - 1, 0)]
            path_stack.append(current_heading)
            current_body = []
            continue
        current_body.append(line)

    flush()
    return CustomerPackPlanner().build_canonical_book(rows, request=record.request)


def _table_to_tagged_text(rows: list[list[str]]) -> str:
    cleaned_rows = [
        [cell.strip() for cell in row]
        for row in rows
        if any(cell.strip() for cell in row)
    ]
    if not cleaned_rows:
        return ""
    rendered_rows = [" | ".join(cell or "-" for cell in row) for row in cleaned_rows]
    return "[TABLE]\n" + "\n".join(rendered_rows) + "\n[/TABLE]"


def _iter_docx_blocks(document: Any):
    if not all((DocxDocumentType, CT_P, CT_Tbl, DocxParagraph, DocxTable)):
        raise RuntimeError("python-docx dependency is unavailable")
    parent = document.element.body
    for child in parent.iterchildren():
        if isinstance(child, CT_P):
            yield DocxParagraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, document)


def _docx_to_structured_text(capture_path: Path) -> str:
    if DocxDocument is None:
        raise RuntimeError("python-docx dependency is unavailable")
    document = DocxDocument(capture_path)
    lines: list[str] = []
    for block in _iter_docx_blocks(document):
        if DocxParagraph is not None and isinstance(block, DocxParagraph):
            text = block.text.strip()
            if not text:
                continue
            style_name = str(getattr(getattr(block, "style", None), "name", "") or "").lower()
            if style_name.startswith("heading"):
                match = re.search(r"(\d+)", style_name)
                level = max(1, min(int(match.group(1)) if match else 1, 6))
                lines.append(f"{'#' * level} {text}")
            else:
                lines.append(text)
            continue
        if DocxTable is not None and isinstance(block, DocxTable):
            table_rows = [
                [cell.text.strip() for cell in row.cells]
                for row in block.rows
            ]
            table_text = _table_to_tagged_text(table_rows)
            if table_text:
                lines.append(table_text)
    return "\n\n".join(line for line in lines if line).strip()


def _pptx_to_structured_text(capture_path: Path) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx dependency is unavailable")
    presentation = Presentation(capture_path)
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        if slide.shapes.title is not None and getattr(slide.shapes.title, "text", "").strip():
            title = slide.shapes.title.text.strip()
        if not title:
            title = f"Slide {index}"
        lines.append(f"# {title}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text and text != title:
                lines.append(text)
            if getattr(shape, "has_table", False):
                table_rows = []
                for row in shape.table.rows:
                    table_rows.append([cell.text.strip() for cell in row.cells])
                table_text = _table_to_tagged_text(table_rows)
                if table_text:
                    lines.append(table_text)
    return "\n\n".join(line for line in lines if line).strip()


def _xlsx_to_structured_text(capture_path: Path) -> str:
    if load_workbook is None:
        raise RuntimeError("openpyxl dependency is unavailable")
    workbook = load_workbook(capture_path, read_only=True, data_only=True)
    try:
        lines: list[str] = []
        for sheet in workbook.worksheets:
            lines.append(f"# {sheet.title}")
            table_rows: list[list[str]] = []
            for row in sheet.iter_rows(values_only=True):
                rendered = ["" if value is None else str(value).strip() for value in row]
                if any(rendered):
                    table_rows.append(rendered)
            table_text = _table_to_tagged_text(table_rows)
            if table_text:
                lines.append(table_text)
        return "\n\n".join(line for line in lines if line).strip()
    finally:
        workbook.close()


def extract_image_markdown_with_docling(path: Path) -> str:
    if DocumentConverter is None:
        raise RuntimeError("docling dependency is unavailable")
    converter = DocumentConverter()
    result = converter.convert(str(path))
    markdown = result.document.export_to_markdown()
    return str(markdown or "").strip()


def _build_docx_canonical_book(record: CustomerPackDraftRecord):
    capture_path = Path(record.capture_artifact_path)
    if not capture_path.exists():
        raise FileNotFoundError(f"captured artifact를 찾을 수 없습니다: {capture_path}")
    normalized = _docx_to_structured_text(capture_path)
    if not normalized:
        raise ValueError("DOCX에서 canonical section을 만들지 못했습니다.")
    return _build_structured_text_canonical_book(record, source_type="md", normalized=normalized)


def _build_pptx_canonical_book(record: CustomerPackDraftRecord):
    capture_path = Path(record.capture_artifact_path)
    if not capture_path.exists():
        raise FileNotFoundError(f"captured artifact를 찾을 수 없습니다: {capture_path}")
    normalized = _pptx_to_structured_text(capture_path)
    if not normalized:
        raise ValueError("PPTX에서 canonical section을 만들지 못했습니다.")
    return _build_structured_text_canonical_book(record, source_type="md", normalized=normalized)


def _build_xlsx_canonical_book(record: CustomerPackDraftRecord):
    capture_path = Path(record.capture_artifact_path)
    if not capture_path.exists():
        raise FileNotFoundError(f"captured artifact를 찾을 수 없습니다: {capture_path}")
    normalized = _xlsx_to_structured_text(capture_path)
    if not normalized:
        raise ValueError("XLSX에서 canonical section을 만들지 못했습니다.")
    return _build_structured_text_canonical_book(record, source_type="md", normalized=normalized)


def _build_image_canonical_book(record: CustomerPackDraftRecord):
    capture_path = Path(record.capture_artifact_path)
    if not capture_path.exists():
        raise FileNotFoundError(f"captured artifact를 찾을 수 없습니다: {capture_path}")
    normalized = _normalize_text_block("md", extract_image_markdown_with_docling(capture_path))
    if not normalized:
        raise ValueError("이미지 OCR에서 canonical section을 만들지 못했습니다.")
    return _build_structured_text_canonical_book(record, source_type="md", normalized=normalized)
