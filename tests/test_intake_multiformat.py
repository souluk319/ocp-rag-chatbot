from __future__ import annotations

import base64
import json
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from docx import Document as DocxDocument
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
TESTS = ROOT / "tests"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))
if str(TESTS) not in sys.path:
    sys.path.insert(0, str(TESTS))

from _support_app_ui import _build_customer_pack_plan, _upload_customer_pack_draft
from play_book_studio.intake import CustomerPackPlanner, DocSourceRequest
from play_book_studio.intake.books.store import CustomerPackDraftStore
from play_book_studio.intake.capture.service import CustomerPackCaptureService
from play_book_studio.intake.normalization.service import CustomerPackNormalizeService


class IntakeMultiformatTests(unittest.TestCase):
    def test_build_customer_pack_plan_accepts_markdown(self) -> None:
        payload = _build_customer_pack_plan(
            {
                "source_type": "md",
                "uri": "/tmp/runbook.md",
                "title": "런북",
            }
        )

        self.assertEqual("md", payload["source_type"])
        self.assertEqual("markdown_text_capture_v1", payload["capture_strategy"])

    def test_upload_customer_pack_draft_preserves_markdown_suffix(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            created = _upload_customer_pack_draft(
                root,
                {
                    "source_type": "md",
                    "uri": "runbook.md",
                    "title": "런북",
                    "file_name": "runbook.md",
                    "content_base64": "IyDrn7zrsJUKCiMjIGV0Y2Qg67Cx7JeF",
                },
            )

        self.assertTrue(str(created["uploaded_file_path"]).endswith(".md"))

    def test_upload_customer_pack_draft_preserves_asciidoc_suffix(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            created = _upload_customer_pack_draft(
                root,
                {
                    "source_type": "asciidoc",
                    "uri": "runbook.adoc",
                    "title": "런북",
                    "file_name": "runbook.adoc",
                    "content_base64": "PSDsmrTsmIEg656c67aBPT0KCl09IGV0Y2Qg67Cx7JeF",
                },
            )

        self.assertTrue(str(created["uploaded_file_path"]).endswith(".adoc"))

    def test_upload_customer_pack_draft_preserves_text_suffix(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            created = _upload_customer_pack_draft(
                root,
                {
                    "source_type": "txt",
                    "uri": "runbook.txt",
                    "title": "런북",
                    "file_name": "runbook.txt",
                    "content_base64": "7Jq07JiBIOyghOywuO2VmOyduA==",
                },
            )

        self.assertTrue(str(created["uploaded_file_path"]).endswith(".txt"))

    def test_markdown_source_uses_text_capture_strategy(self) -> None:
        draft = CustomerPackPlanner().plan(
            DocSourceRequest(
                source_type="md",
                uri="/tmp/runbook.md",
            )
        )

        self.assertEqual("markdown_text_capture_v1", draft.capture_strategy)
        self.assertIn("Markdown", draft.acquisition_step)

    def test_asciidoc_source_uses_text_capture_strategy(self) -> None:
        draft = CustomerPackPlanner().plan(
            DocSourceRequest(
                source_type="asciidoc",
                uri="/tmp/runbook.adoc",
            )
        )

        self.assertEqual("asciidoc_text_capture_v1", draft.capture_strategy)
        self.assertIn("AsciiDoc", draft.acquisition_step)

    def test_text_source_uses_plain_text_capture_strategy(self) -> None:
        draft = CustomerPackPlanner().plan(
            DocSourceRequest(
                source_type="txt",
                uri="/tmp/runbook.txt",
            )
        )

        self.assertEqual("plain_text_capture_v1", draft.capture_strategy)
        self.assertIn("Text", draft.acquisition_step)

    def test_docx_source_uses_binary_capture_strategy(self) -> None:
        draft = CustomerPackPlanner().plan(
            DocSourceRequest(
                source_type="docx",
                uri="/tmp/runbook.docx",
            )
        )

        self.assertEqual("docx_structured_capture_v1", draft.capture_strategy)
        self.assertIn("Word", draft.acquisition_step)

    def test_pptx_source_uses_binary_capture_strategy(self) -> None:
        draft = CustomerPackPlanner().plan(
            DocSourceRequest(
                source_type="pptx",
                uri="/tmp/runbook.pptx",
            )
        )

        self.assertEqual("pptx_slide_capture_v1", draft.capture_strategy)
        self.assertIn("PowerPoint", draft.acquisition_step)

    def test_xlsx_source_uses_binary_capture_strategy(self) -> None:
        draft = CustomerPackPlanner().plan(
            DocSourceRequest(
                source_type="xlsx",
                uri="/tmp/runbook.xlsx",
            )
        )

        self.assertEqual("xlsx_sheet_capture_v1", draft.capture_strategy)
        self.assertIn("Excel", draft.acquisition_step)

    def test_image_source_uses_ocr_capture_strategy(self) -> None:
        draft = CustomerPackPlanner().plan(
            DocSourceRequest(
                source_type="image",
                uri="/tmp/runbook.png",
            )
        )

        self.assertEqual("image_ocr_capture_v1", draft.capture_strategy)
        self.assertIn("Image", draft.acquisition_step)

    def test_capture_service_reads_local_markdown_artifact(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_md = root / "runbook.md"
            source_md.write_text("# 런북\n\n## etcd 백업\n\n```bash\noc debug node/node-0\n```", encoding="utf-8")

            record = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="md", uri=str(source_md), title="런북")
            )

            artifact_path = Path(record.capture_artifact_path)
            self.assertEqual("captured", record.status)
            self.assertTrue(artifact_path.exists())
            self.assertTrue(artifact_path.name.endswith(".md"))

    def test_capture_service_reads_local_asciidoc_artifact(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_adoc = root / "runbook.adoc"
            source_adoc.write_text("= 운영 런북\n\n== etcd 백업\n\n[source,bash]\n----\noc debug node/node-0\n----", encoding="utf-8")

            record = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="asciidoc", uri=str(source_adoc), title="런북")
            )

            artifact_path = Path(record.capture_artifact_path)
            self.assertEqual("captured", record.status)
            self.assertTrue(artifact_path.exists())
            self.assertTrue(artifact_path.name.endswith(".adoc"))

    def test_upload_customer_pack_draft_preserves_docx_suffix(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source = root / "runbook.docx"
            document = DocxDocument()
            document.add_heading("운영 런북", level=1)
            document.save(source)
            created = _upload_customer_pack_draft(
                root,
                {
                    "source_type": "docx",
                    "uri": "runbook.docx",
                    "title": "런북",
                    "file_name": "runbook.docx",
                    "content_base64": base64.b64encode(source.read_bytes()).decode("utf-8"),
                },
            )

        self.assertTrue(str(created["uploaded_file_path"]).endswith(".docx"))

    def test_capture_service_reads_local_docx_artifact(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_docx = root / "runbook.docx"
            document = DocxDocument()
            document.add_heading("운영 런북", level=1)
            document.add_paragraph("etcd 백업 절차")
            document.save(source_docx)

            record = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="docx", uri=str(source_docx), title="런북")
            )

            artifact_path = Path(record.capture_artifact_path)
            self.assertEqual("captured", record.status)
            self.assertTrue(artifact_path.exists())
            self.assertTrue(artifact_path.name.endswith(".docx"))

    def test_capture_service_reads_local_pptx_artifact(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_pptx = root / "runbook.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "운영 런북"
            slide.placeholders[1].text = "etcd 백업 절차"
            presentation.save(source_pptx)

            record = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="pptx", uri=str(source_pptx), title="런북")
            )

            artifact_path = Path(record.capture_artifact_path)
            self.assertEqual("captured", record.status)
            self.assertTrue(artifact_path.exists())
            self.assertTrue(artifact_path.name.endswith(".pptx"))

    def test_capture_service_reads_local_xlsx_artifact(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_xlsx = root / "runbook.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "etcd"
            sheet.append(["절차", "명령어"])
            sheet.append(["백업", "cluster-backup.sh"])
            workbook.save(source_xlsx)

            record = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="xlsx", uri=str(source_xlsx), title="런북")
            )

            artifact_path = Path(record.capture_artifact_path)
            self.assertEqual("captured", record.status)
            self.assertTrue(artifact_path.exists())
            self.assertTrue(artifact_path.name.endswith(".xlsx"))

    def test_capture_service_reads_local_image_artifact(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_png = root / "runbook.png"
            image = Image.new("RGB", (320, 100), color="white")
            draw = ImageDraw.Draw(image)
            draw.text((20, 30), "OpenShift backup restore", fill="black")
            image.save(source_png)

            record = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="image", uri=str(source_png), title="런북")
            )

            artifact_path = Path(record.capture_artifact_path)
            self.assertEqual("captured", record.status)
            self.assertTrue(artifact_path.exists())
            self.assertTrue(artifact_path.name.endswith(".png"))

    def test_normalize_service_builds_sections_from_markdown(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_md = root / "runbook.md"
            source_md.write_text(
                "# 운영 런북\n\n## etcd 백업\n\n```bash\noc debug node/node-0\n```\n\n| 절차 | 명령어 |\n| --- | --- |\n| 백업 | cluster-backup.sh |\n\n## 검증\n\n완료 여부 확인",
                encoding="utf-8",
            )

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="md", uri=str(source_md), title="운영 런북")
            )
            normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
            payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))
            topic_book_path = root / "artifacts" / "customer_packs" / "books" / f"{captured.draft_id}--topic_playbook.json"
            operation_book_path = root / "artifacts" / "customer_packs" / "books" / f"{captured.draft_id}--operation_playbook.json"
            self.assertTrue(topic_book_path.exists())
            self.assertTrue(operation_book_path.exists())

        self.assertEqual("normalized", normalized.status)
        self.assertEqual(2, normalized.normalized_section_count)
        self.assertEqual("md", payload["source_type"])
        self.assertEqual(6, payload["playable_asset_count"])
        self.assertEqual(5, payload["derived_asset_count"])
        self.assertEqual(5, len(payload["derived_assets"]))
        self.assertTrue(any("[CODE" in section["text"] for section in payload["sections"]))
        self.assertTrue(any("[TABLE]" in section["text"] for section in payload["sections"]))
        self.assertEqual(["etcd 백업", "검증"], [section["heading"] for section in payload["sections"]])

    def test_normalize_service_builds_sections_from_asciidoc(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_adoc = root / "runbook.adoc"
            source_adoc.write_text(
                "= 운영 런북\n\n== etcd 백업\n\n[source,bash]\n----\noc debug node/node-0\n----\n\n== 검증\n\n완료 여부 확인",
                encoding="utf-8",
            )

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="asciidoc", uri=str(source_adoc), title="운영 런북")
            )
            normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
            payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertEqual("asciidoc", payload["source_type"])
        self.assertTrue(any("etcd 백업" in section["heading"] for section in payload["sections"]))
        self.assertTrue(
            any("oc debug node/node-0" in section["text"] for section in payload["sections"])
        )

    def test_normalize_service_keeps_more_sections_for_long_derived_playbooks(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_md = root / "long-runbook.md"
            lines = ["# 운영 런북", ""]
            for index in range(1, 21):
                lines.extend(
                    [
                        f"## 단계 {index}",
                        "",
                        f"운영 절차 {index} 확인",
                        "",
                        "```bash",
                        f"oc get pods -n ns-{index}",
                        "```",
                        "",
                    ]
                )
            source_md.write_text("\n".join(lines), encoding="utf-8")

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="md", uri=str(source_md), title="운영 런북")
            )
            normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
            payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        derived_counts = {
            asset["playbook_family"]: int(asset["section_count"])
            for asset in payload["derived_assets"]
        }
        self.assertEqual("normalized", normalized.status)
        self.assertEqual(20, len(payload["sections"]))
        self.assertGreaterEqual(derived_counts["topic_playbook"], 12)
        self.assertGreaterEqual(derived_counts["operation_playbook"], 12)
        self.assertGreaterEqual(derived_counts["troubleshooting_playbook"], 10)
        self.assertGreaterEqual(derived_counts["policy_overlay_book"], 10)
        self.assertGreaterEqual(derived_counts["synthesized_playbook"], 14)

    def test_normalize_service_builds_sections_from_plain_text_numeric_headings(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_txt = root / "runbook.txt"
            source_txt.write_text(
                "1 운영 개요\n\n개요 설명\n\n2.1 etcd 백업\n\noc debug node/node-0\n\n2.2 검증\n\n완료 여부 확인",
                encoding="utf-8",
            )

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="txt", uri=str(source_txt), title="운영 런북")
            )
            normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
            payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        headings = [section["heading"] for section in payload["sections"]]
        self.assertEqual("normalized", normalized.status)
        self.assertEqual("txt", payload["source_type"])
        self.assertTrue(any("2.1 etcd 백업" in heading for heading in headings))
        self.assertTrue(any("2.2 검증" in heading for heading in headings))

    def test_normalize_service_uses_ocr_docling_fallback_for_scan_pdf(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_pdf = root / "scan.pdf"
            source_pdf.write_bytes(b"%PDF-1.4 sample")

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="pdf", uri=str(source_pdf), title="스캔 PDF")
            )

            with patch(
                "play_book_studio.intake.normalization.service.extract_pdf_markdown_with_docling",
                side_effect=ValueError("docling no text"),
            ), patch(
                "play_book_studio.intake.normalization.service.extract_pdf_markdown_with_docling_ocr",
                return_value="## 복구 절차\n\nOCR로 추출된 텍스트",
            ), patch(
                "play_book_studio.intake.normalization.service.extract_pdf_pages",
                side_effect=ValueError("scan pdf"),
            ):
                normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
                payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertEqual("복구 절차", payload["sections"][0]["heading"])

    def test_normalize_service_marks_degraded_pdf_and_blocks_remote_surya_without_approved_boundary(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_pdf = root / "degraded.pdf"
            source_pdf.write_bytes(b"%PDF-1.4 sample")
            (root / ".env").write_text("SURYA_OCR=http://127.0.0.1:8765/ocr\n", encoding="utf-8")

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="pdf", uri=str(source_pdf), title="문제 PDF")
            )

            with patch(
                "play_book_studio.intake.normalization.builders.convert_with_markitdown",
                return_value=(
                    "# 문제 PDF\n\n"
                    "## 온프레미스\n\n"
                    "온프레미스\n\n"
                    "## 설치\n\n"
                    "설치\n\n"
                    "## 네트워크\n\n"
                    "네트워크"
                ),
            ):
                normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
                payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertTrue(normalized.degraded_pdf)
        self.assertIn("too_many_heading_only_sections", normalized.degraded_reason)
        self.assertFalse(normalized.fallback_used)
        self.assertEqual("surya", normalized.fallback_backend)
        self.assertEqual("blocked", normalized.fallback_status)
        self.assertEqual("surya_remote_egress_not_allowed", normalized.fallback_reason)
        self.assertTrue(payload["degraded_pdf"])
        self.assertIn("too_many_heading_only_sections", payload["degraded_reason"])
        self.assertEqual("surya", payload["fallback_backend"])
        self.assertEqual("blocked", payload["fallback_status"])
        self.assertFalse(payload["fallback_used"])

    def test_normalize_service_applies_surya_fallback_for_degraded_pdf_when_boundary_allows_remote_ocr(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_pdf = root / "degraded.pdf"
            source_pdf.write_bytes(b"%PDF-1.4 sample")
            (root / ".env").write_text("SURYA_OCR=http://127.0.0.1:8765/ocr\n", encoding="utf-8")

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="pdf", uri=str(source_pdf), title="문제 PDF")
            )
            captured.tenant_id = "tenant-a"
            captured.workspace_id = "workspace-a"
            captured.access_groups = ("workspace-a", "tenant-a")
            captured.approval_state = "approved"
            captured.publication_state = "draft"
            captured.provider_egress_policy = "surya_remote_ocr"
            captured.redaction_state = "reviewed_safe"
            CustomerPackDraftStore(root).save(captured)

            with patch(
                "play_book_studio.intake.normalization.builders.convert_with_markitdown",
                return_value=(
                    "# 문제 PDF\n\n"
                    "## 온프레미스\n\n"
                    "온프레미스\n\n"
                    "## 설치\n\n"
                    "설치\n\n"
                    "## 네트워크\n\n"
                    "네트워크"
                ),
            ), patch(
                "play_book_studio.intake.normalization.surya_adapter.extract_pdf_markdown_with_surya",
                return_value="# 문제 PDF\n\n## ConfigMap Secret\n\nConfigMap Secret values are synchronized before rollout.",
            ):
                normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
                payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertTrue(normalized.degraded_pdf)
        self.assertTrue(normalized.fallback_used)
        self.assertEqual("surya", normalized.fallback_backend)
        self.assertEqual("applied", normalized.fallback_status)
        self.assertEqual("ConfigMap Secret", payload["sections"][0]["heading"])
        self.assertIn("synchronized before rollout", payload["sections"][0]["text"])

    def test_normalize_service_uses_surya_for_low_confidence_image_when_boundary_allows_remote_ocr(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_png = root / "runbook.png"
            image = Image.new("RGB", (320, 100), color="white")
            draw = ImageDraw.Draw(image)
            draw.text((20, 30), "ConfigMap Secret", fill="black")
            image.save(source_png)
            (root / ".env").write_text("SURYA_OCR=http://127.0.0.1:8765/ocr\n", encoding="utf-8")

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="image", uri=str(source_png), title="런북")
            )
            captured.tenant_id = "tenant-a"
            captured.workspace_id = "workspace-a"
            captured.access_groups = ("workspace-a", "tenant-a")
            captured.approval_state = "approved"
            captured.publication_state = "draft"
            captured.provider_egress_policy = "surya_remote_ocr"
            captured.redaction_state = "reviewed_safe"
            CustomerPackDraftStore(root).save(captured)

            with patch(
                "play_book_studio.intake.normalization.service.extract_image_markdown_with_docling",
                return_value="x",
            ), patch(
                "play_book_studio.intake.normalization.surya_adapter.extract_image_markdown_with_surya",
                return_value="# 런북\n\n## OCR\n\nConfigMap Secret values are synchronized before rollout.",
            ):
                normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
                payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertFalse(normalized.degraded_pdf)
        self.assertTrue(normalized.fallback_used)
        self.assertEqual("surya", normalized.fallback_backend)
        self.assertEqual("applied", normalized.fallback_status)
        self.assertIn("ConfigMap Secret values are synchronized", payload["sections"][0]["text"])

    def test_normalize_service_does_not_flag_healthy_pdf_as_degraded(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_pdf = root / "healthy.pdf"
            source_pdf.write_bytes(b"%PDF-1.4 sample")

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="pdf", uri=str(source_pdf), title="정상 PDF")
            )

            with patch(
                "play_book_studio.intake.normalization.builders.convert_with_markitdown",
                return_value=(
                    "# 정상 PDF\n\n"
                    "## 백업 절차\n\n"
                    "cluster-backup.sh 실행 후 상태를 확인합니다.\n\n"
                    "## 검증\n\n"
                    "완료 상태와 관련 로그를 함께 검토합니다."
                ),
            ):
                normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
                payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertFalse(normalized.degraded_pdf)
        self.assertEqual("", normalized.degraded_reason)
        self.assertFalse(normalized.fallback_used)
        self.assertEqual("", normalized.fallback_backend)
        self.assertEqual("not_needed", normalized.fallback_status)
        self.assertFalse(payload["degraded_pdf"])
        self.assertEqual("", payload["degraded_reason"])
        self.assertEqual("ready", payload["quality_status"])
        self.assertEqual("not_needed", payload["fallback_status"])

    def test_normalize_service_cleans_pdf_markitdown_page_markers_and_tableish_headings(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_pdf = root / "storage.pdf"
            source_pdf.write_bytes(b"%PDF-1.4 sample")

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="pdf", uri=str(source_pdf), title="02. 스토리지(03.19)")
            )

            with patch(
                "play_book_studio.intake.normalization.builders.convert_with_markitdown",
                return_value=(
                    "# 02. 스토리지(03.19)\n\n"
                    "스토리지 10\n\n"
                    "1 NFS ( | --- | --- | --- | )\n"
                    "스토리지 공유 개요\n\n"
                    "## PV 생성\n\n"
                    "apiVersion: v1\n"
                    "kind: PersistentVolume\n\n"
                    "스토리지 11\n\n"
                    "## PVC 확인\n\n"
                    "PVC\x00정상 확인"
                ),
            ):
                normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
                payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        texts = [section["text"] for section in payload["sections"]]
        headings = [section["heading"] for section in payload["sections"]]
        self.assertEqual("normalized", normalized.status)
        self.assertEqual(["02. 스토리지(03.19)", "PV 생성", "PVC 확인"], headings)
        self.assertFalse(any("|" in heading for heading in headings))
        self.assertFalse(any("\x00" in text for text in texts))
        self.assertFalse(any("스토리지 10" in text or "스토리지 11" in text for text in texts))

    def test_normalize_service_repairs_sparse_pdf_markdown_tables_into_compact_rows(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_pdf = root / "configmap-secret.pdf"
            source_pdf.write_bytes(b"%PDF-1.4 sample")

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="pdf", uri=str(source_pdf), title="ConfigMap Secret")
            )

            with patch(
                "play_book_studio.intake.normalization.builders.convert_with_markitdown",
                return_value=(
                    "# ConfigMap Secret\n\n"
                    "## ConfigMap Secret 비교\n\n"
                    "| ConfigMap: | | | 일반 | 설정 데이터 | 저장 | | | |\n"
                    "| --- | --- | --- | --- | --- | --- | --- | --- | --- |\n"
                    "| Secret: | | 민감한 | 정보 | 저장 | | | | |\n"
                ),
            ):
                normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
                payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        section_text = payload["sections"][0]["text"]
        self.assertIn('[TABLE header="false"]', section_text)
        self.assertIn("ConfigMap | 일반 설정 데이터 저장", section_text)
        self.assertIn("Secret | 민감한 정보 저장", section_text)
        self.assertNotIn("| - | - |", section_text)

    def test_normalize_service_does_not_split_numeric_table_rows_into_headings(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_pdf = root / "secret-yaml.pdf"
            source_pdf.write_bytes(b"%PDF-1.4 sample")

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="pdf", uri=str(source_pdf), title="Secret YAML")
            )

            with patch(
                "play_book_studio.intake.normalization.builders.convert_with_markitdown",
                return_value=(
                    "# Secret YAML\n\n"
                    "## Secret 생성\n\n"
                    "| 단계 | 설명 |\n"
                    "| --- | --- |\n"
                    "| 7.1 YAML | Secret 생성 |\n"
                    "| 7.2 CLI | 명령 예시 |\n\n"
                    "다음 단계 설명\n\n"
                    "## 검증\n\n"
                    "완료 여부를 확인합니다.\n"
                ),
            ):
                normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
                payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        headings = [section["heading"] for section in payload["sections"]]
        self.assertEqual(["Secret 생성", "검증"], headings)
        self.assertIn("7.1 YAML | Secret 생성", payload["sections"][0]["text"])
        self.assertIn("[/TABLE]\n\n다음 단계 설명", payload["sections"][0]["text"])

    def test_docx_source_is_honestly_rejected(self) -> None:
        payload = _build_customer_pack_plan({"source_type": "docx", "uri": "/tmp/demo.docx"})
        self.assertEqual("docx", payload["source_type"])
        self.assertEqual("docx_structured_capture_v1", payload["capture_strategy"])

    def test_normalize_service_builds_sections_from_docx(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_docx = root / "runbook.docx"
            document = DocxDocument()
            document.add_heading("운영 런북", level=1)
            document.add_heading("etcd 백업", level=2)
            document.add_paragraph("oc debug node/node-0")
            table = document.add_table(rows=2, cols=2)
            table.rows[0].cells[0].text = "절차"
            table.rows[0].cells[1].text = "명령어"
            table.rows[1].cells[0].text = "백업"
            table.rows[1].cells[1].text = "cluster-backup.sh"
            document.save(source_docx)

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="docx", uri=str(source_docx), title="운영 런북")
            )
            normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
            payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertEqual("docx", payload["source_type"])
        self.assertTrue(any("etcd 백업" in section["heading"] for section in payload["sections"]))
        self.assertTrue(any("[TABLE]" in section["text"] for section in payload["sections"]))
        self.assertTrue(any("MarkItDown" in note for note in payload.get("normalization_notes", [])))

    def test_normalize_service_builds_sections_from_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_pptx = root / "runbook.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "etcd 백업"
            slide.placeholders[1].text = "cluster-backup.sh 실행"
            presentation.save(source_pptx)

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="pptx", uri=str(source_pptx), title="운영 런북")
            )
            normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
            payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertEqual("pptx", payload["source_type"])
        self.assertTrue(any("etcd 백업" in section["heading"] for section in payload["sections"]))
        self.assertTrue(any("cluster-backup.sh" in section["text"] for section in payload["sections"]))
        self.assertTrue(any("MarkItDown" in note for note in payload.get("normalization_notes", [])))

    def test_normalize_service_builds_sections_from_xlsx(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_xlsx = root / "runbook.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "backup"
            sheet.append(["절차", "명령어"])
            sheet.append(["백업", "cluster-backup.sh"])
            workbook.save(source_xlsx)

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="xlsx", uri=str(source_xlsx), title="운영 런북")
            )
            normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
            payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertEqual("xlsx", payload["source_type"])
        self.assertTrue(any("backup" in section["heading"].lower() for section in payload["sections"]))
        self.assertTrue(any("[TABLE]" in section["text"] for section in payload["sections"]))
        self.assertTrue(any("MarkItDown" in note for note in payload.get("normalization_notes", [])))

    def test_normalize_service_builds_sections_from_image_ocr(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_png = root / "runbook.png"
            image = Image.new("RGB", (320, 100), color="white")
            draw = ImageDraw.Draw(image)
            draw.text((20, 30), "OpenShift backup restore", fill="black")
            image.save(source_png)

            captured = CustomerPackCaptureService(root).capture(
                request=DocSourceRequest(source_type="image", uri=str(source_png), title="OCR 런북")
            )

            with patch(
                "play_book_studio.intake.normalization.service.extract_image_markdown_with_docling",
                return_value="## OCR 절차\n\ncluster-backup.sh",
            ):
                normalized = CustomerPackNormalizeService(root).normalize(draft_id=captured.draft_id)
                payload = json.loads(Path(normalized.canonical_book_path).read_text(encoding="utf-8"))

        self.assertEqual("normalized", normalized.status)
        self.assertEqual("image", payload["source_type"])
        self.assertTrue(any("OCR 절차" in section["heading"] for section in payload["sections"]))


if __name__ == "__main__":
    unittest.main()
